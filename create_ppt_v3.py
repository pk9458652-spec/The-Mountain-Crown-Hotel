from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# --- STYLE CONSTANTS ---
BG_COLOR = RGBColor(15, 23, 42)      # Deep Midnight
ACCENT_BLUE = RGBColor(56, 189, 248) # Sky Blue
ACCENT_GREEN = RGBColor(34, 197, 94) # Emerald
TEXT_WHITE = RGBColor(255, 255, 255)
TEXT_GRAY = RGBColor(148, 163, 184)

def apply_attractive_background(slide):
    # Add a full-slide rectangle for a consistent deep-theme background
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5)
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    line = shape.line
    line.fill.background() # No border

    # Add a subtle accent line at the bottom
    line_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(7.3), Inches(10), Inches(0.2)
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = ACCENT_BLUE
    line_shape.line.fill.background()

def add_styled_heading(slide, text, y_pos=0.5, color=ACCENT_BLUE):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.name = 'Impact'
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    return txBox

def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_attractive_background(slide)

    # Main Title with Shadow Effect (offset text)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "THE MOUNTAIN CROWN"
    p.font.bold = True
    p.font.size = Pt(66)
    p.font.name = 'Impact'
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Accent Line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(3.4), Inches(5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(1.0))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "NEXT-GEN HOTEL MANAGEMENT ECOSYSTEM"
    p2.font.size = Pt(24)
    p2.font.name = 'Century Gothic'
    p2.font.color.rgb = ACCENT_BLUE
    p2.alignment = PP_ALIGN.CENTER

def add_affiliation_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_attractive_background(slide)

    add_styled_heading(slide, "Academic Affiliation", color=TEXT_WHITE)

    # College Name in a 'Glass' box
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.8), Inches(8), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(30, 41, 59)
    shape.line.color.rgb = ACCENT_BLUE

    tf = shape.text_frame
    tf.vertical_anchor = 1 # Top
    p = tf.paragraphs[0]
    p.text = "THAKUR JAGDEV CHAND MEMORIAL DEGREE COLLEGE"
    p.font.bold = True
    p.font.size = Pt(26)
    p.font.name = 'Segoe UI Semibold'
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "SUJANPUR TIRA"
    p2.font.size = Pt(20)
    p2.alignment = PP_ALIGN.CENTER

    # Team Members
    tx3 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(4), Inches(2))
    tf3 = tx3.text_frame
    tf3.text = "DEVELOPED BY"
    tf3.paragraphs[0].font.bold = True
    tf3.paragraphs[0].font.color.rgb = ACCENT_BLUE
    tf3.paragraphs[0].font.size = Pt(20)
    
    for name in ["Pardeep Kumar (23BCA045)", "Anshul Koundal (23BCA028)"]:
        p = tf3.add_paragraph()
        p.text = "• " + name
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_WHITE

    # University
    tx4 = slide.shapes.add_textbox(Inches(5.5), Inches(3.8), Inches(4), Inches(2))
    tf4 = tx4.text_frame
    tf4.text = "SUPERVISED BY"
    tf4.paragraphs[0].font.bold = True
    tf4.paragraphs[0].font.color.rgb = ACCENT_BLUE
    tf4.paragraphs[0].font.size = Pt(20)
    
    p = tf4.add_paragraph()
    p.text = "Faculty of BCA Dept."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_WHITE
    p = tf4.add_paragraph()
    p.text = "Himachal Pradesh University"
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_GRAY

def add_toc_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_attractive_background(slide)
    add_styled_heading(slide, "Project Roadmap", y_pos(0.4))

    contents = [
        "Core Introduction", "Tech Stack Overview", "Development Phases",
        "Implementation Flow", "Visual Showcase", "Security Challenges", "Future Horizons"
    ]

    for i, item in enumerate(contents):
        y = 1.8 + (i * 0.7)
        # Bullet Shape
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(y), Inches(7), Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(30, 41, 59)
        shape.line.color.rgb = ACCENT_BLUE
        
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = f"{i+1}. {item}"
        p.font.size = Pt(18)
        p.font.name = 'Impact'
        p.font.color.rgb = TEXT_WHITE
        p.alignment = PP_ALIGN.LEFT

def y_pos(val): return Inches(val)

def add_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_attractive_background(slide)
    add_styled_heading(slide, "The Hospitality Vision", color=TEXT_WHITE)

    # Styled Text Blocks
    for i, (title, content, color) in enumerate([
        ("THE CHALLENGE", "Manual ledgers and third-party commissions (15-25%) drain hotel profits and cause synchronization errors.", RGBColor(239, 68, 68)),
        ("OUR SOLUTION", "A bespoke, zero-commission digital portal providing absolute control over bookings and inventory.", ACCENT_GREEN)
    ]):
        y = 2.0 + (i * 2.2)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(y), Inches(8), Inches(1.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(30, 41, 59)
        shape.line.width = Pt(2)
        shape.line.color.rgb = color
        
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(24)
        p.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = content
        p2.font.size = Pt(18)
        p2.font.color.rgb = TEXT_WHITE

def add_styled_list_slide(prs, title_str, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_attractive_background(slide)
    add_styled_heading(slide, title_str)

    for i, (head, desc) in enumerate(items):
        y = 1.8 + (i * 1.5) # Increased spacing
        # Dot
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(y+0.1), Inches(0.2), Inches(0.2))
        oval.fill.solid()
        oval.fill.fore_color.rgb = ACCENT_BLUE
        oval.line.fill.background()
        
        tx = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(8), Inches(1))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = head.upper()
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_WHITE
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(15)
        p2.font.color.rgb = TEXT_GRAY

def add_visual_slide(prs, title, path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_attractive_background(slide)
    add_styled_heading(slide, title)

    if os.path.exists(path):
        # Large Image Frame
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.8), Inches(8), Inches(4.5))
        shape.fill.background()
        shape.line.color.rgb = ACCENT_BLUE
        shape.line.width = Pt(3)
        
        # Actual Image (Center it)
        slide.shapes.add_picture(path, Inches(1.1), Inches(1.9), width=Inches(7.8))

def add_future_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_attractive_background(slide)
    add_styled_heading(slide, "Future Roadmap", color=ACCENT_GREEN)

    points = [
        ("PAYMENT GATEWAY", "Automated UPI & Card processing."),
        ("AI PRICING", "Smart demand-based rate adjustment."),
        ("MOBILE APP", "React Native extension for iOS/Android."),
        ("IoT LOCKS", "Direct room access via digital keys.")
    ]
    
    for i, (h, d) in enumerate(points):
        y = 2.0 + (i * 1.2)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(y), Inches(8), Inches(1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(30, 41, 59)
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = h + ": " + d
        p.font.bold = True
        p.font.size = Pt(22)
        p.font.color.rgb = TEXT_WHITE
        p.alignment = PP_ALIGN.CENTER

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. Title
    add_title_slide(prs)
    # 2. Affiliation
    add_affiliation_slide(prs)
    # 3. TOC
    add_toc_slide(prs)
    # 4. Intro
    add_intro_slide(prs)
    
    # 5. Tech Stack & Features
    tech_items = [
        ("FLASK & SQLITE", "Secure Python core mapping to robust SQLite data storage."),
        ("GUEST PORTAL", "Check booking status securely with 10-digit phone validation."),
        ("ADMIN PANEL", "Invisible logo access, room floor allotment & entry deletion.")
    ]
    add_styled_list_slide(prs, "Architecture & Key Features", tech_items)

    # 6. Progress
    progress_items = [
        ("PHASE 1", "Comprehensive market research & gap analysis."),
        ("PHASE 2", "Schema design & API endpoint mapping."),
        ("PHASE 3", "Full-stack integration & rigorous testing.")
    ]
    add_styled_list_slide(prs, "Strategic Development", progress_items)

    # 7. (Removed ER/DFD as requested) - Skip or replace with something else?
    # User said "remove the er and dfd diagram from the 7 slide". I'll just skip it.

    # 8. Visual Showcases (Multiple Slides)
    add_visual_slide(prs, "Homepage Overview", "home_screen.jpg")
    add_visual_slide(prs, "Luxury Rooms Section", "rooms_screen.jpg")
    add_visual_slide(prs, "Photo Gallery Experience", "gallery_screen.jpg")
    add_visual_slide(prs, "Secure Admin Dashboard", "admin_screen.jpg")

    # 9. Future
    add_future_slide(prs)

    # 10. Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_attractive_background(slide)
    
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    p = tx.text_frame.paragraphs[0]
    p.text = "THANK YOU"
    p.font.bold = True
    p.font.size = Pt(96)
    p.font.name = 'Impact'
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
    p2 = tx2.text_frame.paragraphs[0]
    p2.text = "Questions & Discussions Invited"
    p2.font.size = Pt(24)
    p2.font.color.rgb = ACCENT_BLUE
    p2.alignment = PP_ALIGN.CENTER

    output_file = "Mountain_Crown_Premium_Presentation.pptx"
    prs.save(output_file)
    print(f"Created: {output_file}")

if __name__ == "__main__":
    main()
