from pptx import Presentation

def get_ppt_structure(file_path):
    prs = Presentation(file_path)
    for i, slide in enumerate(prs.slides):
        title = slide.shapes.title.text if slide.shapes.title else "No Title"
        print(f"Slide {i+1}: {title}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text != title:
                print(f"  - {shape.text[:100].replace('\n', ' ')}")

if __name__ == "__main__":
    get_ppt_structure("Smart-Solar-Dustbin-with-Waste-Management-and-SMS-Alert-System (3) (1)-1.pptx")
