from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def set_slide_background(slide, color_rgb):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb

def add_title_slide(prs):
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, RGBColor(15, 23, 42)) # Dark Blue/Black

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "THE MOUNTAIN CROWN"
    p.font.bold = True
    p.font.size = Pt(60)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1.0))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Hotel Management and Luxury Booking System"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(148, 163, 184)
    p2.alignment = PP_ALIGN.CENTER

    # Tagline
    txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Innovative hospitality management using Flask and IoT-inspired web tech"
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(56, 189, 248)
    p3.alignment = PP_ALIGN.CENTER

def add_affiliation_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, RGBColor(255, 255, 255))

    # College Name
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    p = tx.text_frame.paragraphs[0]
    p.text = "THAKUR JAGDEV CHAND MEMORIAL DEGREE COLLEGE\nSUJANPUR TIRA"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(30, 41, 59)
    p.alignment = PP_ALIGN.CENTER

    # Dept
    tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(0.5))
    p2 = tx2.text_frame.paragraphs[0]
    p2.text = "Department of Computer Applications (BCA)"
    p2.font.size = Pt(22)
    p2.alignment = PP_ALIGN.CENTER

    # Team
    tx3 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(4), Inches(2))
    tf3 = tx3.text_frame
    tf3.text = "Submitted By:"
    tf3.paragraphs[0].font.bold = True
    p = tf3.add_paragraph()
    p.text = "Pardeep Kumar (23BCA045)"
    p = tf3.add_paragraph()
    p.text = "Anshul Koundal (23BCA028)"

    # Supervisor
    tx4 = slide.shapes.add_textbox(Inches(5), Inches(3.5), Inches(4), Inches(2))
    tf4 = tx4.text_frame
    tf4.text = "Submitted To:"
    tf4.paragraphs[0].font.bold = True
    p = tf4.add_paragraph()
    p.text = "Department Faculty Members"
    p = tf4.add_paragraph()
    p.text = "Himachal Pradesh University"

def add_toc_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, RGBColor(248, 250, 252))

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title.text_frame.text = "Table of Contents"
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(36)

    contents = [
        ("01", "Introduction", "The Hospitality Challenge & Our Digital Solution"),
        ("02", "Technology Stack", "Full-Stack Specifications & Database Design"),
        ("03", "Project Progress", "From Requirement Analysis to Deployment"),
        ("04", "System Design", "DFD, ER Diagram and Core Logic"),
        ("05", "Live Demo & UI", "Visualizing The Mountain Crown Experience"),
        ("06", "Challenges Faced", "Security & Synchronization Solutions"),
        ("07", "Future Outcomes", "Scalability, Payments, and AI enhancements")
    ]

    for i, (num, head, desc) in enumerate(contents):
        y = 1.5 + (i * 0.8)
        # Number
        tx1 = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(0.5), Inches(0.5))
        tx1.text_frame.text = num
        tx1.text_frame.paragraphs[0].font.bold = True
        tx1.text_frame.paragraphs[0].font.color.rgb = RGBColor(56, 189, 248)
        
        # Heading
        tx2 = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(3), Inches(0.5))
        tx2.text_frame.text = head
        tx2.text_frame.paragraphs[0].font.bold = True
        
        # Description
        tx3 = slide.shapes.add_textbox(Inches(4.5), Inches(y), Inches(5), Inches(0.5))
        tx3.text_frame.text = desc
        tx3.text_frame.paragraphs[0].font.size = Pt(14)
        tx3.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 116, 139)

def add_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, RGBColor(255, 255, 255))

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title.text_frame.text = "Introduction: The Hospitality Challenge & Solution"
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(28)

    # Left Column: The Challenge
    tx1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.2), Inches(3))
    tf1 = tx1.text_frame
    p = tf1.paragraphs[0]
    p.text = "The Traditional Challenge"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(220, 38, 38)
    
    p = tf1.add_paragraph()
    p.text = "Hotels often rely on manual paper-based ledgers or expensive third-party agencies (OTAs) charging 15-25% commissions. This leads to double-bookings and revenue loss."
    p.font.size = Pt(16)

    # Right Column: Our Innovation
    tx2 = slide.shapes.add_textbox(Inches(5.3), Inches(1.5), Inches(4.2), Inches(3))
    tf2 = tx2.text_frame
    p = tf2.paragraphs[0]
    p.text = "Our Digital Innovation"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(22, 163, 74)
    
    p = tf2.add_paragraph()
    p.text = "A bespoke, zero-commission web portal with a secure admin dashboard. It automates room availability, simplifies guest booking, and puts full control back in the hotel owner's hands."
    p.font.size = Pt(16)

def add_tech_stack_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title.text_frame.text = "Technology Stack: Powering The Mountain Crown"
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(28)

    techs = [
        ("Flask (Python 3.10)", "The robust 'brain' of the system. Handles routing, API logic, and secure authentication."),
        ("SQLite Database", "ACID-compliant storage for rooms, bookings, and customer data. High reliability, zero maintenance."),
        ("Vanilla JS & CSS3", "Lightning-fast frontend with modern glassmorphism effects. Optimized for SEO and mobile speed."),
        ("RESTful API Architecture", "Seamless communication between guest interface and backend booking engine."),
        ("Security: SQL Injection Defense", "Implemented parameterized queries and secure master passcodes for admin access.")
    ]

    for i, (head, desc) in enumerate(techs):
        y = 1.5 + (i * 1.1)
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(1))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = head
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 112, 192)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(71, 85, 105)

def add_progress_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title.text_frame.text = "Project Progress: From Concept to Deployment"
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(28)

    steps = [
        ("1", "Phase 1: Research", "Analyzed existing booking systems and identified gaps in commission-based models."),
        ("2", "Phase 2: Design", "Created UI/UX wireframes and mapped out the SQLite relational schema."),
        ("3", "Phase 3: Development", "Built the Flask API and integrated Vanilla JS components for a smooth user experience."),
        ("4", "Phase 4: Testing", "Rigorous testing of date-overlap logic and security patches for the admin dashboard."),
        ("5", "Phase 5: Deployment", "Final production build optimized and deployed for local/cloud accessibility.")
    ]

    for i, (num, head, desc) in enumerate(steps):
        y = 1.5 + (i * 1.1)
        # Circle/Number
        tx1 = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(0.5), Inches(0.5))
        tx1.text_frame.text = num
        tx1.text_frame.paragraphs[0].font.bold = True
        tx1.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Head & Desc
        tx2 = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(8), Inches(1))
        tf2 = tx2.text_frame
        p = tf2.paragraphs[0]
        p.text = head
        p.font.bold = True
        p.font.size = Pt(18)
        
        p = tf2.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(71, 85, 105)

def add_design_slide(prs, title_str, img_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title.text_frame.text = title_str
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(28)

    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8))

def add_problems_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title.text_frame.text = "Challenges & Solutions Implemented"
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(28)

    challenges = [
        ("Complexity of Date Logic", "Problem: Booking overlaps caused scheduling conflicts.\nSolution: Implemented a robust date-aware availability engine in the backend."),
        ("Admin Data Security", "Problem: Managing dashboard access without heavy frameworks.\nSolution: Designed a secure, passcode-guarded route with session-based logic."),
        ("Asset Loading Speed", "Problem: HD images slowed down the initial page load.\nSolution: Optimized image delivery and used modern CSS for pre-loading effects.")
    ]

    for i, (head, desc) in enumerate(challenges):
        y = 1.5 + (i * 1.5)
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(1.5))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = head
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(220, 38, 38)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(30, 41, 59)

def add_future_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, RGBColor(15, 23, 42))

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title.text_frame.text = "Future Scope & Enhancements"
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(28)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    futures = [
        ("UPI & Payment Integration", "Enabling direct bookings with automated payment confirmations via UPI/Stripe."),
        ("AI-Driven Dynamic Pricing", "Using AI to adjust room rates based on seasonal demand and occupancy trends."),
        ("Mobile App (React Native)", "Extending the ecosystem to Android and iOS for easier guest access."),
        ("IoT Room Automation", "Connecting the portal to smart locks and energy management systems.")
    ]

    for i, (head, desc) in enumerate(futures):
        y = 1.5 + (i * 1.2)
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(1))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = head
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(56, 189, 248)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(148, 163, 184)

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_affiliation_slide(prs)
    add_toc_slide(prs)
    add_intro_slide(prs)
    add_tech_stack_slide(prs)
    add_progress_slide(prs)
    add_design_slide(prs, "System Design: DFD & ER Diagram", "er_screen.jpg")
    add_design_slide(prs, "Live Experience: UI & Admin Panel", "admin_screen.jpg")
    add_problems_slide(prs)
    add_future_slide(prs)

    # Closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, RGBColor(15, 23, 42))
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1.5))
    p = tx.text_frame.paragraphs[0]
    p.text = "THANK YOU!"
    p.font.bold = True
    p.font.size = Pt(80)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    output_file = "Mountain_Crown_Final_Presentation.pptx"
    prs.save(output_file)
    print(f"Created: {output_file}")

if __name__ == "__main__":
    main()
