from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def add_title_slide(prs):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "THE MOUNTAIN CROWN\nHOTEL MANAGEMENT SYSTEM"
    subtitle.text = ("Project Presentation\n\nSubmitted By:\nPARDEEP KUMAR (23BCA045)\nANSHUL KOUNDAL (23BCA028)\n\nHIMACHAL PRADESH UNIVERSITY")

    # Styling
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 112, 192) # Blue

def add_slide(prs, title_str, points, image_path=None):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_str
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = points[0] if points else ""
    
    for point in points[1:]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0

    if image_path and os.path.exists(image_path):
        # Add image to the side or bottom
        # Content usually takes the left, we can shrink it or put image on right
        body_shape.width = Inches(5.5)
        slide.shapes.add_picture(image_path, Inches(6.0), Inches(1.5), height=Inches(4.0))

def main():
    prs = Presentation()

    # 1. Title Slide
    add_title_slide(prs)

    # 2. Introduction
    add_slide(prs, "Introduction", [
        "Modern digital portal for luxury hotel management.",
        "Goal: Automate reservations and enhance guest experience.",
        "Built for responsiveness, security, and speed.",
        "Integrates high-definition visuals with a robust booking engine."
    ])

    # 3. Problem Statement
    add_slide(prs, "Problem Statement", [
        "Heavy reliance on manual paper-based ledgers.",
        "High commission fees (15-25%) from third-party travel agencies.",
        "Vulnerability to double-bookings and human oversight.",
        "Lack of centralized data for management to track revenue."
    ])

    # 4. Objectives
    add_slide(prs, "Project Objectives", [
        "Develop an independent, zero-commission booking platform.",
        "Implement a date-aware conflict prevention engine.",
        "Design a secure, hidden Admin Dashboard for staff.",
        "Ensure mobile-first responsiveness and high performance."
    ])

    # 5. Proposed System
    add_slide(prs, "Proposed System Features", [
        "Modern Vanilla JS Frontend: Lightning-fast load times.",
        "Flask Backend: Robust and secure routing logic.",
        "SQLite Database: ACID-compliant data integrity.",
        "Admin Portal: Real-time revenue and inventory management."
    ], "home_screen.jpg")

    # 6. System Requirements
    add_slide(prs, "System Requirements", [
        "Software: Python 3.10+, Flask, SQLite3, HTML5, CSS3, JS.",
        "Hardware: Intel i3/i5 Processor, 4GB+ RAM.",
        "Tools: VS Code, Git, Chrome DevTools.",
        "Server: Local Flask server / Any WSGI-compliant server."
    ])

    # 7. Data Flow Diagram (DFD)
    add_slide(prs, "System Design: DFD", [
        "Level 0: High-level interaction between Guest, Website, and Admin.",
        "Level 1: Detail of Booking processes and DB queries.",
        "Visualizes how data flows into the storage layers."
    ], "dfd_screen.jpg")

    # 8. ER Diagram
    add_slide(prs, "System Design: ER Diagram", [
        "Entities: Guest, Room, Booking, Review.",
        "Relationships: One-to-Many mapping for Bookings.",
        "Ensures relational data integrity across tables."
    ], "er_screen.jpg")

    # 9. Implementation: Frontend
    add_slide(prs, "Implementation: UI & UX", [
        "Responsive Rooms Section with live pricing.",
        "Dynamic Gallery fetching images in HD.",
        "Bespoke 'Pay at Hotel' checkout flow.",
        "Modern theme with glassmorphism effects."
    ], "rooms_screen.jpg")

    # 10. Implementation: Backend & Security
    add_slide(prs, "Implementation: Backend & Security", [
        "Parameterized Queries to prevent SQL Injection.",
        "Date-overlap algorithm for reservation safety.",
        "Secure passcode protection for Admin Access.",
        "RESTful API design for modular frontend connectivity."
    ], "admin_screen.jpg")

    # 11. Testing & Results
    add_slide(prs, "Testing & Results", [
        "Unit Testing of API endpoints.",
        "Integration Testing between UI and Flask.",
        "Performance: Under 1.5s page load speed.",
        "Zero critical bugs identified in final UAT."
    ])

    # 12. Conclusion & Future Scope
    add_slide(prs, "Conclusion & Future Scope", [
        "Successfully developed an end-to-end hotel portal.",
        "Future: Payment Gateway integration (Stripe/UPI).",
        "Future: Automatic SMS / Email confirmations.",
        "Future: AI-driven pricing based on demand seasonality."
    ])

    prs.save('Mountain_Crown_Project_Presentation.pptx')
    print("Presentation created successfully: Mountain_Crown_Project_Presentation.pptx")

if __name__ == "__main__":
    main()
